import csv
from tkinter import END, LEFT, filedialog, simpledialog
import customtkinter as ctk
from tkinter import Listbox as TkListbox  

import subprocess
from pathlib import Path

#manipular objetos do powerpoint
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor

ctk.set_appearance_mode('dark')
app = ctk.CTk()
app.title('Gerador de Certificados - PPG-Fil')
app.geometry('900x900')

nome_congresso = ''

caminhoPPTX = ''

def selecionar_pasta():
    pasta = filedialog.askdirectory()
    if pasta:
        entry_path_folder.delete(0, ctk.END)
        entry_path_folder.insert(0, pasta)

def pptx_para_pdf(pptx_path):
    pptx_path = Path(pptx_path).resolve()
    pasta_destino = pptx_path.parent

    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(pasta_destino),
        str(pptx_path)
    ])

def extrair_listbox_para_lista(listbox, destino):
    # Limpa a lista de destino
    destino.clear()

    # Itera sobre todos os itens da Listbox
    for i in range(listbox.size()):
        valor = listbox.get(i)
        destino.append(valor)

def extrair_nm_titles(arquivo_csv: str, temp_author, temp_titles: list):
    with open(arquivo_csv, newline='', encoding='utf-8') as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            # Garante que a linha tenha pelo menos duas colunas
            if len(row) >= 2:
                temp_author.append(row[0].strip())
                temp_titles.append(row[1].strip())


def selecionar_arquivo():
    caminho = filedialog.askopenfilename(
        title="Selecione um arquivo CSV",
        filetypes=[("CSV files", "*.csv"), ("Todos os arquivos", "*.*")]
    )
    if caminho:
        entry_path_csv.delete(0, ctk.END)
        entry_path_csv.insert(0, caminho)

def selecionar_arquivo_pptx():
    caminho = filedialog.askopenfilename(
        title="Selecione um arquivo PPTX",
        filetypes=[("PPTX Files", "*.pptx"), ("Todos os arquivos", "*.*")]
    )
    if caminho:
        entry_path_pptx.delete(0, ctk.END)
        entry_path_pptx.insert(0, caminho)

def adicionar_autor():
    temp = simpledialog.askstring("Entrada de texto", "Digite o nome da caixa de texto:")
    listaAuthor.insert(END, temp)

def remover_autor():
    selecao = listaAuthor.curselection()
    if selecao:
        listaAuthor.delete(selecao)

def adicionar_title_presentation():
    temp = simpledialog.askstring("Entrada de texto", "Digite o nome da caixa de texto:")
    listaTitle.insert(END, temp)

def remover_title_presentation():
    selecao = listaTitle.curselection()
    if selecao:
        listaTitle.delete(selecao)

def modificar_conteudo_apresentacao_pptx(arquivo_pptx, caminho_pdf, novo_conteudo_nome, novo_conteudo_titulo):
    presentation = Presentation(arquivo_pptx)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.name in cx_texto_autor:
                    shape.text = novo_conteudo_nome
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    shape.text_frame.paragraphs[0].font.size = Pt(17)
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    shape.text_frame.paragraphs[0].font.name = "Lato"
                elif shape.name in cx_texto_titulo:
                    shape.text = novo_conteudo_titulo
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    shape.text_frame.paragraphs[0].font.name = "Lato"
                    shape.text_frame.paragraphs[0].font.size = Pt(13)
                

    caminho = f'{caminho_pdf}{nome_congresso} - {novo_conteudo_nome}'
    caminho_pptx_temp = f'{caminho}.pptx'
    
    presentation.save(caminho_pptx_temp)
    pptx_para_pdf(caminho_pptx_temp)

def gerarPDF():
    caminhoPPTX = entry_path_pptx.get()
    caminho_destino = f'{entry_path_folder.get()}/'
    extrair_nm_titles(entry_path_csv.get(), nm_authors, nm_titles)
    extrair_listbox_para_lista(listaAuthor, cx_texto_autor)
    extrair_listbox_para_lista(listaTitle,cx_texto_titulo)
    
    for i in range(len(nm_authors)):
        modificar_conteudo_apresentacao_pptx(caminhoPPTX, caminho_destino,nm_authors[i], nm_titles[i])

###### INÍCIO - criação de campos: Arquivo CSV, Modelo em PPTX. ######
label_caminho_csv = ctk.CTkLabel(app,text='Caminho do Arquivo CSV')
label_caminho_csv.pack(pady=10)


entry_path_csv = ctk.CTkEntry(app, width=400, placeholder_text="Caminho do arquivo")
entry_path_csv.pack(pady=10)

botao_arquivo = ctk.CTkButton(app, text="Selecionar Arquivo", command=selecionar_arquivo)
botao_arquivo.pack()

label_caminho_csv = ctk.CTkLabel(app,text='Caminho do Modelo de Certificado')
label_caminho_csv.pack(pady=10)


entry_path_pptx = ctk.CTkEntry(app, width=400, placeholder_text="Caminho do arquivo")
entry_path_pptx.pack(pady=10)
###### FIM - criação de campos: Arquivo CSV, Modelo em PPTX. ######

# Botão para selecionar arquivo
botao_arquivo = ctk.CTkButton(app, text="Selecionar Arquivo", command=selecionar_arquivo_pptx)
botao_arquivo.pack()


######## INICÍO - Frame horizontal para as duas listas lado a lado ########
label_boxs = ctk.CTkLabel(app, text="Caixas de texto para substituição")
label_boxs.pack(pady=10)
frame_listas = ctk.CTkFrame(app)
frame_listas.pack(pady=10)

# Lista de autores(as)
frame_autores = ctk.CTkFrame(frame_listas)
frame_autores.pack(side=LEFT, padx=20)

label_autores = ctk.CTkLabel(frame_autores, text="Caixas: Nome dos(as) autores(as)")
label_autores.pack(pady=5)

listaAuthor = TkListbox(frame_autores, height=10, width=40)
listaAuthor.pack()

frame_botoes_autores = ctk.CTkFrame(frame_autores, fg_color="transparent")
frame_botoes_autores.pack(pady=5)

botao_inserir_autor = ctk.CTkButton(frame_botoes_autores, text="Inserir", command=adicionar_autor)
botao_inserir_autor.pack(side=LEFT, padx=5)

botao_remover_autor = ctk.CTkButton(frame_botoes_autores, text="Remover", command=remover_autor)
botao_remover_autor.pack(side=LEFT, padx=5)

# Lista de títulos das apresentações
frame_titulos = ctk.CTkFrame(frame_listas)
frame_titulos.pack(side=LEFT, padx=20)

label_titulos = ctk.CTkLabel(frame_titulos, text="Caixas: Título da apresentação")
label_titulos.pack(pady=5)

listaTitle = TkListbox(frame_titulos, height=10, width=40)
listaTitle.pack()

frame_botoes_titulos = ctk.CTkFrame(frame_titulos, fg_color="transparent")
frame_botoes_titulos.pack(pady=5)

botao_inserir_titulo = ctk.CTkButton(frame_botoes_titulos, text="Inserir", command=adicionar_title_presentation)
botao_inserir_titulo.pack(side=LEFT, padx=5)

botao_remover_titulo = ctk.CTkButton(frame_botoes_titulos, text="Remover", command=remover_title_presentation)
botao_remover_titulo.pack(side=LEFT, padx=5)
######## FIM - Frame horizontal para as duas listas lado a lado ########

######## Caminho da pasta de destino dos arquivos em pdf ###########
label_pasta_destino = ctk.CTkLabel(app,text='Caminho da pasta destino para os arquivos em PDF')
label_pasta_destino.pack(pady=10)


entry_path_folder = ctk.CTkEntry(app, width=400, placeholder_text="Caminho da pasta")
entry_path_folder.pack(pady=10)

######## Botão para selecionar a pasta de destino ########
botao = ctk.CTkButton(app, text="Selecionar Pasta", command=selecionar_pasta)
botao.pack(pady=10)
######## Caminho da pasta de destino dos arquivos em pdf ###########


#Listas que recebem os valores das caixas de texto a serem substituídas.
cx_texto_autor = []
cx_texto_titulo = []

#Nome do congresso (ou simposósio, workshop, etc.) que irá no nome do arquivo final. E.g.: '13th Principia Symposium - Nome do Autor.pdf' 
nome_congresso = simpledialog.askstring("Entrada de texto", "Digite o nome do Congresso:")

# Botão -- Gerar PDF -- 
botao_gerar_pdf = ctk.CTkButton(app,text='Gerar Certificados em PDF', command = gerarPDF)
botao_gerar_pdf.pack()

######### Listas para armazenar os nomes de autores(as) e das apresentações #########
nm_authors = []
nm_titles = []

#Abertura do ambiente gráfico.
app.mainloop()