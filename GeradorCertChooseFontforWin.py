import csv
from tkinter import END, LEFT, filedialog, simpledialog, StringVar, colorchooser
import customtkinter as ctk
from tkinter import Listbox as TkListbox  
import comtypes.client
from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor

ctk.set_appearance_mode('dark')
app = ctk.CTk()
app.title('Gerador de Certificados - PPG-Fil')
app.geometry('1200x1000')

# Font customization variables
author_font = {
    'name': 'Arial',
    'size': 17,
    'color': (255, 255, 255),
    'alignment': PP_ALIGN.CENTER
}

title_font = {
    'name': 'Arial',
    'size': 13,
    'color': (255, 255, 255),
    'alignment': PP_ALIGN.JUSTIFY
}

nome_congresso = ''
caminhoPPTX = ''

def selecionar_pasta():
    pasta = filedialog.askdirectory()
    if pasta:
        entry_path_folder.delete(0, ctk.END)
        entry_path_folder.insert(0, pasta)

def pptx_para_pdf(pptx_path):
    pptx_path = Path(pptx_path).resolve()
    pasta_destino = pptx_path.parent
    nome_arquivo_pdf = pptx_path.with_suffix(".pdf")

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    try:
        apresentacao = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        apresentacao.SaveAs(str(nome_arquivo_pdf), FileFormat=32)  # 32 = PDF
        apresentacao.Close()
    finally:
        powerpoint.Quit()

def extrair_listbox_para_lista(listbox, destino):
    destino.clear()
    for i in range(listbox.size()):
        valor = listbox.get(i)
        destino.append(valor)

def extrair_nm_titles(arquivo_csv: str, temp_author, temp_titles: list):
    with open(arquivo_csv, newline='', encoding='utf-8') as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            if len(row) >= 2:
                temp_author.append(row[0].strip())
                temp_titles.append(row[1].strip())

def selecionar_arquivo():
    caminho = filedialog.askopenfilename(
        title="Selecione um arquivo CSV",
        filetypes=[("CSV files", "*.csv"), ("Todos os arquivos", "*.*")]
    )
    if caminho:
        entry_path_csv.delete(0, ctk.END)
        entry_path_csv.insert(0, caminho)

def selecionar_arquivo_pptx():
    caminho = filedialog.askopenfilename(
        title="Selecione um arquivo PPTX",
        filetypes=[("PPTX Files", "*.pptx"), ("Todos os arquivos", "*.*")]
    )
    if caminho:
        entry_path_pptx.delete(0, ctk.END)
        entry_path_pptx.insert(0, caminho)

def adicionar_autor():
    temp = simpledialog.askstring("Entrada de texto", "Digite o nome da caixa de texto:")
    if temp:  # Only insert if user didn't cancel
        listaAuthor.insert(END, temp)

def remover_autor():
    selecao = listaAuthor.curselection()
    if selecao:
        listaAuthor.delete(selecao)

def adicionar_title_presentation():
    temp = simpledialog.askstring("Entrada de texto", "Digite o nome da caixa de texto:")
    if temp:  # Only insert if user didn't cancel
        listaTitle.insert(END, temp)

def remover_title_presentation():
    selecao = listaTitle.curselection()
    if selecao:
        listaTitle.delete(selecao)

# New font customization functions
def choose_author_font_name():
    font = simpledialog.askstring("Font Name", "Enter font name for author:", initialvalue=author_font['name'])
    if font:
        author_font['name'] = font
        author_font_name_label.configure(text=f"Font: {font}")

def choose_author_font_size():
    size = simpledialog.askinteger("Font Size", "Enter font size for author:", initialvalue=author_font['size'])
    if size:
        author_font['size'] = size
        author_font_size_label.configure(text=f"Size: {size}pt")



def choose_title_font_name():
    font = simpledialog.askstring("Font Name", "Enter font name for title:", initialvalue=title_font['name'])
    if font:
        title_font['name'] = font
        title_font_name_label.configure(text=f"Font: {font}")

def choose_title_font_size():
    size = simpledialog.askinteger("Font Size", "Enter font size for title:", initialvalue=title_font['size'])
    if size:
        title_font['size'] = size
        title_font_size_label.configure(text=f"Size: {size}pt")

def choose_author_font_color():
    color = colorchooser.askcolor(title="Choose author font color", initialcolor=(author_font['color'][0], author_font['color'][1], author_font['color'][2]))
    if color[1]:
        author_font['color'] = (int(color[0][0]), int(color[0][1]), int(color[0][2]))
        author_font_color_label.configure(fg_color=color[1], text=f"Color: RGB({author_font['color'][0]}, {author_font['color'][1]}, {author_font['color'][2]})")

def choose_title_font_color():
    color = colorchooser.askcolor(title="Choose title font color", initialcolor=(title_font['color'][0], title_font['color'][1], title_font['color'][2]))
    if color[1]:
        title_font['color'] = (int(color[0][0]), int(color[0][1]), int(color[0][2]))
        title_font_color_label.configure(fg_color=color[1], text=f"Color: RGB({title_font['color'][0]}, {title_font['color'][1]}, {title_font['color'][2]})")

def modificar_conteudo_apresentacao_pptx(arquivo_pptx, caminho_pdf, novo_conteudo_nome, novo_conteudo_titulo):
    presentation = Presentation(arquivo_pptx)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.name in cx_texto_autor:
                    shape.text = novo_conteudo_nome
                    shape.text_frame.paragraphs[0].alignment = author_font['alignment']
                    shape.text_frame.paragraphs[0].font.size = Pt(author_font['size'])
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*author_font['color'])
                    shape.text_frame.paragraphs[0].font.name = author_font['name']
                elif shape.name in cx_texto_titulo:
                    shape.text = novo_conteudo_titulo
                    shape.text_frame.paragraphs[0].alignment = title_font['alignment']
                    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_font['color'])
                    shape.text_frame.paragraphs[0].font.name = title_font['name']
                    shape.text_frame.paragraphs[0].font.size = Pt(title_font['size'])

    caminho = f'{caminho_pdf}{nome_congresso} - {novo_conteudo_nome}'
    caminho_pptx_temp = f'{caminho}.pptx'
    
    presentation.save(caminho_pptx_temp)
    pptx_para_pdf(caminho_pptx_temp)

def gerarPDF():
    caminhoPPTX = entry_path_pptx.get()
    caminho_destino = f'{entry_path_folder.get()}/'
    extrair_nm_titles(entry_path_csv.get(), nm_authors, nm_titles)
    extrair_listbox_para_lista(listaAuthor, cx_texto_autor)
    extrair_listbox_para_lista(listaTitle,cx_texto_titulo)
    
    for i in range(len(nm_authors)):
        modificar_conteudo_apresentacao_pptx(caminhoPPTX, caminho_destino,nm_authors[i], nm_titles[i])

# Opções de alinhamento das fontes -- Dict types!
alignment_options = {
    "Left": PP_ALIGN.LEFT,
    "Center": PP_ALIGN.CENTER,
    "Right": PP_ALIGN.RIGHT,
    "Justify": PP_ALIGN.JUSTIFY
}

def set_author_alignment(choice):
    author_font['alignment'] = alignment_options[choice]

def set_title_alignment(choice):
    title_font['alignment'] = alignment_options[choice]

###### UI Elements ######
# File selection section
label_caminho_csv = ctk.CTkLabel(app, text='Caminho arquivo CSV:')
label_caminho_csv.pack(pady=5)

entry_path_csv = ctk.CTkEntry(app, width=400, placeholder_text="Caminho arquivo CSV")
entry_path_csv.pack(pady=5)

botao_arquivo = ctk.CTkButton(app, text="Selecionar arquivo CSV", command=selecionar_arquivo)
botao_arquivo.pack(pady=5)

label_caminho_csv = ctk.CTkLabel(app, text='Modelo de Certificado (PPTX):')
label_caminho_csv.pack(pady=5)

entry_path_pptx = ctk.CTkEntry(app, width=400, placeholder_text="Caminho arquivo PPTX")
entry_path_pptx.pack(pady=5)

botao_arquivo = ctk.CTkButton(app, text="Selecionar arquivo PPTX", command=selecionar_arquivo_pptx)
botao_arquivo.pack(pady=5)

# Font customization frames
font_customization_frame = ctk.CTkFrame(app)
font_customization_frame.pack(pady=5, padx=5)

# Author font customization
author_font_frame = ctk.CTkFrame(font_customization_frame)
author_font_frame.pack(side="left", padx=5, pady=5)

ctk.CTkLabel(author_font_frame, text="Config Fonte dos Autores").pack(pady=5)

# Font name
ctk.CTkButton(author_font_frame, text="Definir nome da fonte", command=choose_author_font_name).pack(pady=5)
author_font_name_label = ctk.CTkLabel(author_font_frame, text=f"Font: {author_font['name']}")
author_font_name_label.pack()

# Font size
ctk.CTkButton(author_font_frame, text="Definir tamanho da fonte", command=choose_author_font_size).pack(pady=5)
author_font_size_label = ctk.CTkLabel(author_font_frame, text=f"Size: {author_font['size']}pt")
author_font_size_label.pack()

# Font color
ctk.CTkButton(author_font_frame, text="Definir cor da fonte", command=choose_author_font_color).pack(pady=5)
author_font_color_label = ctk.CTkLabel(author_font_frame, width=30, height=5, 
                                     text=f"Color: RGB({author_font['color'][0]}, {author_font['color'][1]}, {author_font['color'][2]})")
author_font_color_label.pack()

# Font alignment
ctk.CTkLabel(author_font_frame, text="Alinhamento:").pack(pady=5)
author_font_alignment = ctk.CTkOptionMenu(author_font_frame, values=list(alignment_options.keys()), command=set_author_alignment)
author_font_alignment.pack()
author_font_alignment.set("Center")

# Title font customization
title_font_frame = ctk.CTkFrame(font_customization_frame)
title_font_frame.pack(side="right", padx=20, pady=10)

ctk.CTkLabel(title_font_frame, text="Configurações da fonte do título").pack(pady=5)

# Font name
ctk.CTkButton(title_font_frame, text="Definir nome da fonte", command=choose_title_font_name).pack(pady=5)
title_font_name_label = ctk.CTkLabel(title_font_frame, text=f"Fonte: {title_font['name']}")
title_font_name_label.pack()

# Font size
ctk.CTkButton(title_font_frame, text="Definir tamanho da fonte", command=choose_title_font_size).pack(pady=5)
title_font_size_label = ctk.CTkLabel(title_font_frame, text=f"Tamanho: {title_font['size']}pt")
title_font_size_label.pack()

# Font color
ctk.CTkButton(title_font_frame, text="Definir cor da fonte", command=choose_title_font_color).pack(pady=5)
title_font_color_label = ctk.CTkLabel(title_font_frame, width=30, height=5, 
                                     text=f"Cor: RGB({title_font['color'][0]}, {title_font['color'][1]}, {title_font['color'][2]})")
title_font_color_label.pack()

# Font alignment
ctk.CTkLabel(title_font_frame, text="Alinhamento:").pack(pady=5)
title_font_alignment = ctk.CTkOptionMenu(title_font_frame, values=list(alignment_options.keys()), command=set_title_alignment)
title_font_alignment.pack()
title_font_alignment.set("Justify")

# Boxes section
label_boxs = ctk.CTkLabel(app, text="Caixas de texto a serem substituídas")
label_boxs.pack(pady=5)

frame_listas = ctk.CTkFrame(app)
frame_listas.pack(pady=5)

# Author list
frame_autores = ctk.CTkFrame(frame_listas)
frame_autores.pack(side=LEFT, padx=5)

label_autores = ctk.CTkLabel(frame_autores, text="Caixa de texto de nomes dos(as) autores(as)")
label_autores.pack(pady=5)

listaAuthor = TkListbox(frame_autores, height=10, width=40)
listaAuthor.pack()

frame_botoes_autores = ctk.CTkFrame(frame_autores, fg_color="transparent")
frame_botoes_autores.pack(pady=5)

botao_inserir_autor = ctk.CTkButton(frame_botoes_autores, text="Inserir", command=adicionar_autor)
botao_inserir_autor.pack(side=LEFT, padx=5)

botao_remover_autor = ctk.CTkButton(frame_botoes_autores, text="Remover", command=remover_autor)
botao_remover_autor.pack(side=LEFT, padx=5)

# Title list
frame_titulos = ctk.CTkFrame(frame_listas)
frame_titulos.pack(side=LEFT, padx=5)

label_titulos = ctk.CTkLabel(frame_titulos, text="Caixa de texto de títulos de apresentações")
label_titulos.pack(pady=5)

listaTitle = TkListbox(frame_titulos, height=10, width=40)
listaTitle.pack()

frame_botoes_titulos = ctk.CTkFrame(frame_titulos, fg_color="transparent")
frame_botoes_titulos.pack(pady=5)

botao_inserir_titulo = ctk.CTkButton(frame_botoes_titulos, text="Inserir", command=adicionar_title_presentation)
botao_inserir_titulo.pack(side=LEFT, padx=5)

botao_remover_titulo = ctk.CTkButton(frame_botoes_titulos, text="Remover", command=remover_title_presentation)
botao_remover_titulo.pack(side=LEFT, padx=5)

# Destination folder section
label_pasta_destino = ctk.CTkLabel(app, text='Pasta de destino para arquivos .PDF')
label_pasta_destino.pack(pady=5)

entry_path_folder = ctk.CTkEntry(app, width=400, placeholder_text="Caminho pasta de destino")
entry_path_folder.pack(pady=5)

frame_botoes_finais = ctk.CTkFrame(app)
frame_botoes_finais.pack(padx=5)

botao = ctk.CTkButton(frame_botoes_finais, text="Selecionar pasta", command=selecionar_pasta)
botao.pack(side=LEFT,padx=5)

# List to store text box names
cx_texto_autor = []
cx_texto_titulo = []


# Generate PDF button
botao_gerar_pdf = ctk.CTkButton(frame_botoes_finais, text='Gerar Certificados em PDF', command=gerarPDF)
botao_gerar_pdf.pack(side=LEFT,padx=5)

# Get congress name
nome_congresso = simpledialog.askstring("Nome do congresso", "Insira o nome do congresso:")

# Lists for authors and titles
nm_authors = []
nm_titles = []

app.mainloop()
